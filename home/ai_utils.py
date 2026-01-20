import os
from groq import Groq
import PyPDF2
from pptx import Presentation
from pathlib import Path
from dotenv import load_dotenv
import json

# ==========================================
# 🔑 CONFIGURATION
# ==========================================
current_file = Path(__file__).resolve()
project_dir = current_file.parent.parent
search_paths = [project_dir / '.env', project_dir.parent / '.env', current_file.parent / '.env']

env_path = None
for path in search_paths:
    if path.exists():
        env_path = path
        break

if env_path:
    load_dotenv(env_path)

api_key = os.getenv("GROQ_API_KEY")

# ==========================================
# 📄 FILE PROCESSING
# ==========================================
def extract_text_from_pdf(file_obj):
    try:
        reader = PyPDF2.PdfReader(file_obj)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"PDF Error: {e}")
        return None

def extract_text_from_ppt(file_obj):
    try:
        prs = Presentation(file_obj)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"PPT Error: {e}")
        return None

# ==========================================
# 🧠 AI FUNCTIONS
# ==========================================

def generate_notes(uploaded_file):
    """
    Returns a TUPLE: (HTML_Summary, Raw_Text)
    We need the Raw_Text to perform the 'Deep Dives' later.
    """
    if not api_key:
        return "❌ Error: GROQ_API_KEY not found.", ""

    # 1. Extract Text
    filename = uploaded_file.name.lower()
    text = ""
    
    if filename.endswith('.pdf'):
        text = extract_text_from_pdf(uploaded_file)
    elif filename.endswith('.pptx') or filename.endswith('.ppt'):
        text = extract_text_from_ppt(uploaded_file)
    else:
        return "❌ Error: Unsupported file format.", ""

    if not text or len(text) < 50:
        return "❌ Error: Could not extract text.", ""

    # 2. Call Groq API (Initial Summary)
    try:
        client = Groq(api_key=api_key)
        truncated_text = text[:50000] # 50k char limit

        completion = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a precise academic analyst. Output HTML only."
                },
                {
                    "role": "user",
                    "content": f"""
                    Analyze this text. Structure the output as follows:

                    1. Break into SUBTOPICS.
                    2. For each subtopic, use an <h3> tag with the specific Topic Name.
                    3. Under the <h3>, provide a bulleted summary <ul>.
                    4. End with a 🏁 Final Summary paragraph.

                    Example Format:
                    <h3>📌 [Topic Name]</h3>
                    <ul><li>Point 1</li></ul>
                    <hr>

                    TEXT:
                    {truncated_text}
                    """
                }
            ],
            model="llama-3.3-70b-versatile",
        )
        
        # RETURN BOTH!
        return completion.choices[0].message.content, truncated_text

    except Exception as e:
        return f"❌ AI Error: {str(e)}", ""


def generate_deep_dive(topic, full_text):
    """
    Generates a detailed explanation for a SPECIFIC topic using the full context.
    """
    if not api_key:
        return "Error: API Key missing."

    try:
        client = Groq(api_key=api_key)
        
        completion = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a strict academic tutor. You provide exhaustive, detailed explanations."
                },
                {
                    "role": "user",
                    "content": f"""
                    CONTEXT:
                    {full_text}
                    
                    TASK:
                    The student has requested a "Deep Dive" on the specific topic: "{topic}".
                    
                    INSTRUCTIONS:
                    1. Reread the Context and find EVERYTHING related to "{topic}".
                    2. Explain it in extreme detail. Do not leave out any nuance, formula, or figure mentioned in the text.
                    3. If there are steps, list them. If there is a table described, format it as HTML.
                    4. Use simple HTML (<p>, <ul>, <strong>, <table>).
                    
                    OUTPUT:
                    Provide ONLY the detailed explanation.
                    """
                }
            ],
            model="llama-3.3-70b-versatile",
        )
        return completion.choices[0].message.content

    except Exception as e:
        return f"Error: {str(e)}"

# home/ai_utils.py
def generate_quiz(text_content, num_questions=10):
    """
    Generates a JSON quiz from the provided text.
    """
    if not api_key:
        return None

    try:
        client = Groq(api_key=api_key)
        
        # Limit text to avoid token limits (approx 15k chars is safe for Llama-70b context)
        truncated_text = text_content[:20000]

        prompt = f"""
        Generate {num_questions} multiple-choice questions based strictly on the text provided below.
        
        OUTPUT FORMAT (JSON ONLY):
        [
            {{
                "question": "Question text here?",
                "options": ["Option A", "Option B", "Option C", "Option D"],
                "answer": "Option B"
            }},
            ...
        ]
        
        RULES:
        1. Return ONLY valid JSON. No markdown formatting, no ```json tags.
        2. Ensure exactly {num_questions} questions.
        3. Make options plausible.
        
        TEXT CONTENT:
        {truncated_text}
        """

        completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": "You are a teacher creating a quiz. Output raw JSON only."},
                {"role": "user", "content": prompt}
            ],
            model="llama-3.3-70b-versatile",
            temperature=0.5, # Lower temperature for factual accuracy
        )
        
        # Clean up potential markdown formatting from AI
        raw_content = completion.choices[0].message.content
        raw_content = raw_content.replace('```json', '').replace('```', '').strip()
        
        return json.loads(raw_content)

    except Exception as e:
        print(f"Quiz Generation Error: {e}")
        return None